from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

BG_COLOR = RGBColor(0x13, 0x15, 0x1a)
ACCENT_COLOR = RGBColor(0x00, 0xd4, 0xff)
WHITE = RGBColor(255, 255, 255)

def set_dark_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()
    return bg

slide_title = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_background(slide_title)
tx1 = slide_title.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(2))
tf1 = tx1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "Inteligencia Artificial II"
p.font.size = Pt(80)
p.font.bold = True
p.font.color.rgb = ACCENT_COLOR
p.alignment = PP_ALIGN.CENTER

tx2 = slide_title.shapes.add_textbox(Inches(2), Inches(5), Inches(12), Inches(2))
tf2 = tx2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Búsqueda y Optimización en Almacenes\nExposición de Algoritmos"
p2.font.size = Pt(36)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER

def create_concept_slide(title_text, img_path, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(32)
    p2.font.color.rgb = WHITE
    
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(3), Inches(2.6), width=Inches(10))

def create_technical_slide(title_text, plot_path, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    if plot_path and os.path.exists(plot_path):
        slide.shapes.add_picture(plot_path, Inches(1), Inches(2), height=Inches(6))
        
    txBox2 = slide.shapes.add_textbox(Inches(8), Inches(2.5), Inches(7.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for bp in bullet_points:
        p2 = tf2.add_paragraph()
        p2.text = "• " + bp
        p2.font.size = Pt(32)
        p2.font.color.rgb = WHITE
        p2.space_after = Pt(24)

img_concept_warehouse = r"C:\Users\Juan\.gemini\antigravity\brain\72f26122-54e1-4b55-a0fc-598721940235\warehouse_robots_1773840661175.png"
img_concept_multi = r"C:\Users\Juan\.gemini\antigravity\brain\72f26122-54e1-4b55-a0fc-598721940235\abstract_multiagent_1773840752966.png"
img_concept_gen = r"C:\Users\Juan\.gemini\antigravity\brain\72f26122-54e1-4b55-a0fc-598721940235\abstract_genetic_1773840817067.png"

create_concept_slide("1. Búsqueda A* (A-Star)", img_concept_warehouse, "Concepto: Navegación esquivando obstáculos fijos de forma óptima.")
create_technical_slide("A*: Implementación y Resultado", "img_astar.png", ["Algoritmo guiado por heurística (Distancia Manhattan).", "Garantiza encontrar el camino más corto.", "Base para algoritmos más complejos."])

create_concept_slide("2. Sistema Multi-Agente", img_concept_multi, "Concepto: Coordinación cooperativa sin colisiones.")
create_technical_slide("Multi-Agente: Resultado", "img_doble.png", ["Múltiples robots operando simultáneamente.", "Resolución de colisiones re-evaluando rutas dinámicamente.", "Alta escalabilidad para flotas de picking."])

create_concept_slide("3. Temple Simulado (Simulated Annealing)", None, "Concepto: Optimización de la ruta de picking (metaheurística).")
create_technical_slide("Temple Simulado: Convergencia", "img_temple.png", ["Encuentra el recorrido más rápido entre múltiples estanterías.", "Evita estancarse en óptimos locales aceptando probabilísticamente peores soluciones al inicio."])

create_concept_slide("4. Algoritmo Genético", img_concept_gen, "Concepto: Evolución de las mejores distribuciones del almacén.")
create_technical_slide("Genético: Ubicación por Frecuencia", "img_genetico.png", ["Intercambia productos buscando que los más pedidos queden cerca.", "Cruce y mutación aseguran diversidad de configuraciones.", "Reduce drásticamente el tiempo promedio de picking."])

slide_final = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_background(slide_final)
txf = slide_final.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(2))
tff = txf.text_frame
pf = tff.paragraphs[0]
pf.text = "¡Gracias por su atención!"
pf.font.size = Pt(72)
pf.font.bold = True
pf.font.color.rgb = ACCENT_COLOR
pf.alignment = PP_ALIGN.CENTER

prs.save("Presentacion_Super_Pro.pptx")
print("Presentación creada exitosamente: Presentacion_Super_Pro.pptx")
