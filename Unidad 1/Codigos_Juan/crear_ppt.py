from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Inteligencia Artificial II"
subtitle.text = "Búsqueda y Optimización en Almacenes\n(Enfoques y Resolución)"

def add_slide(title_text, img_path, bullet_points):
    layout = prs.slide_layouts[5] # blank slide with title
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = title_text
    
    if os.path.exists(img_path):
        try:
            # Position image on left, size 5 inches width
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.8), width=Inches(5))
        except Exception as e:
            print(f"Error cargando {img_path}: {e}")
    else:
        print(f"Advertencia: No se encontró la imagen {img_path}")
        
    # Text box on the right
    txBox = slide.shapes.add_textbox(Inches(5.8), Inches(2), Inches(3.8), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(20)

add_slide("Búsqueda A* (A-Star)", "img_astar.png", ["Búsqueda del camino más corto.", "Detección y evasión de obstáculos fijos (estanterías)."])
add_slide("Sistema Multi-Agente", "img_doble.png", ["Múltiples robots operando simultáneamente.", "Resolución de colisiones en tiempo real.", "Recálculo rápido de ruteo eficiente."])
add_slide("Temple Simulado", "img_temple.png", ["Cálculo de secuencia óptima (ruta de picking).", "Evita óptimos locales mediante aceptación probabilística.", "Optimiza distancia total del recorrido."])
add_slide("Algoritmo Genético", "img_genetico.png", ["Distribución inteligente del almacén.", "Productos más frecuentes al inicio (menor distancia).", "Evolución por selección natural (cruzamiento y mutación)."])

prs.save("Presentacion_IA.pptx")
print("Presentación creada exitosamente: Presentacion_IA.pptx")
